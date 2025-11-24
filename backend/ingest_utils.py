# backend/ingest_utils.py
import os
import json

# PyMuPDF (fitz) is optional for PDF support; import safely
try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None

# python-docx for DOCX support
try:
    from docx import Document
except Exception:
    Document = None

# python-pptx for PPTX support
try:
    from pptx import Presentation
except Exception:
    Presentation = None

# Pillow and pytesseract for image OCR
try:
    from PIL import Image
    import pytesseract
    
    # Configure Tesseract path for Windows
    # Common installation paths
    import platform
    if platform.system() == "Windows":
        possible_paths = [
            r"C:\Program Files\Tesseract-OCR\tesseract.exe",
            r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
            r"C:\Users\{}\AppData\Local\Programs\Tesseract-OCR\tesseract.exe".format(os.environ.get('USERNAME', ''))
        ]
        for path in possible_paths:
            if os.path.exists(path):
                pytesseract.pytesseract.tesseract_cmd = path
                break
except Exception:
    Image = None
    pytesseract = None


def extract_text_from_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def extract_text_from_pdf(path: str) -> str:
    if fitz is None:
        raise RuntimeError("PyMuPDF (fitz) not installed. pip install pymupdf")
    doc = fitz.open(path)
    parts = []
    for page in doc:
        parts.append(page.get_text())
    return "\n".join(parts)


def extract_text_from_docx(path: str) -> str:
    """Extract text from DOCX files."""
    if Document is None:
        raise RuntimeError("python-docx not installed. pip install python-docx")
    
    doc = Document(path)
    parts = []
    
    # Extract paragraphs
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            parts.append(paragraph.text)
    
    # Extract tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)
    
    return "\n".join(parts)


def extract_text_from_pptx(path: str) -> str:
    """Extract text from PPTX files."""
    if Presentation is None:
        raise RuntimeError("python-pptx not installed. pip install python-pptx")
    
    prs = Presentation(path)
    parts = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        parts.append(f"\n--- Slide {slide_num} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    
    return "\n".join(parts)


def extract_text_from_image(path: str) -> str:
    """Extract text from images using OCR (Tesseract)."""
    if Image is None or pytesseract is None:
        raise RuntimeError("Pillow and pytesseract not installed. pip install Pillow pytesseract")
    
    try:
        img = Image.open(path)
        text = pytesseract.image_to_string(img)
        return text
    except Exception as e:
        print(f"OCR error for {path}: {e}")
        return f"[Image file - OCR failed: {str(e)}]"


def extract_text(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    
    # Text-based formats
    if ext in [".txt", ".md", ".csv"]:
        return extract_text_from_txt(path)
    
    # PDF
    if ext in [".pdf"]:
        return extract_text_from_pdf(path)
    
    # DOCX
    if ext in [".docx"]:
        return extract_text_from_docx(path)
    
    # PPTX
    if ext in [".pptx"]:
        return extract_text_from_pptx(path)
    
    # Images (OCR)
    if ext in [".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".gif"]:
        return extract_text_from_image(path)
    
    # Fallback: try reading as text
    try:
        return extract_text_from_txt(path)
    except Exception:
        return f"[Unsupported file format: {ext}]"


def save_preview(content_text: str, original_path: str, preview_dir: str = None, chars: int = 2000):
    """
    Saves a small preview JSON for quick inspection.
    Returns the preview file path.
    """
    if preview_dir is None:
        preview_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "uploads")
    os.makedirs(preview_dir, exist_ok=True)
    base = os.path.basename(original_path)
    preview_name = f"{base}.preview.json"
    preview_path = os.path.join(preview_dir, preview_name)
    preview = {
        "source_file": os.path.basename(original_path),
        "preview": content_text[:chars],
        "length": len(content_text),
    }
    with open(preview_path, "w", encoding="utf-8") as f:
        json.dump(preview, f, ensure_ascii=False, indent=2)
    return preview_path
